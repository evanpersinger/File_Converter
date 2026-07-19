"""Convert PowerPoint (.pptx) files to Markdown.

For each .pptx in input/, walks every slide with python-pptx, adds a `## Slide N`
heading, promotes large-font text to subheadings, and writes the text to output/.
"""

import os
import glob
from pptx import Presentation

# Get the directory where this script is located
script_dir = os.path.dirname(os.path.abspath(__file__))

# Folders
input_folder = os.path.join(script_dir, 'input')
output_folder = os.path.join(script_dir, 'output')


def extract_text_from_pptx(pptx_file):
    """Extract text from PowerPoint file and format as markdown"""
    prs = Presentation(pptx_file)
    markdown_content = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        markdown_content.append(f"## Slide {slide_num}\n")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                # Check if this is likely a title
                if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                    first_para = shape.text_frame.paragraphs[0]
                    if first_para.runs and first_para.runs[0].font.size:
                        if first_para.runs[0].font.size > 200000:
                            markdown_content.append(f"### {text}\n")
                            continue
                
                markdown_content.append(f"{text}\n")
        
        markdown_content.append("\n---\n\n")
    
    return ''.join(markdown_content)


def convert_pptx_to_markdown() -> str:
    """Convert all PPTX files in the input folder to Markdown files in the output folder.

    Returns:
        A summary of what was converted, suitable for showing to a caller.
    """

    # Create output folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Find all PPTX files
    pptx_files = glob.glob(os.path.join(input_folder, '*.pptx'))

    if not pptx_files:
        # If there are only markdown files present, notify the user
        md_present = glob.glob(os.path.join(input_folder, '*.md'))
        if md_present:
            return "That file is already in markdown format"
        return "No PPTX files found in input folder"

    print(f"Found {len(pptx_files)} PPTX files to convert")

    converted = []
    errors = []

    for pptx_file in pptx_files:
        try:
            # Get filename without extension
            filename = os.path.splitext(os.path.basename(pptx_file))[0]
            md_file = os.path.join(output_folder, f"{filename}.md")

            # Extract text and format as markdown
            markdown_content = extract_text_from_pptx(pptx_file)

            # Write to markdown file
            with open(md_file, 'w', encoding='utf-8') as f:
                f.write(markdown_content)

            print(f"Converted: {os.path.basename(pptx_file)} -> {filename}.md")
            converted.append(f"{filename}.md")

        except Exception as e:
            print(f"Error converting {pptx_file}: {str(e)}")
            errors.append(f"{os.path.basename(pptx_file)}: {e}")

    if not converted:
        return f"No files converted. {len(errors)} failed: {'; '.join(errors)}"

    summary = f"Converted {len(converted)} file(s) to output/: {', '.join(converted)}"
    if errors:
        summary += f". {len(errors)} failed: {'; '.join(errors)}"
    return summary


if __name__ == "__main__":
    print(convert_pptx_to_markdown())
